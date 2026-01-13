from pptx import Presentation
import os

def extract_pptx_text(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return

    prs = Presentation(pptx_path)
    
    print(f"--- Slides in {os.path.basename(pptx_path)} ---")
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}:")
        
        # Get title
        if slide.shapes.title:
            print(f"  Title: {slide.shapes.title.text}")
        
        # Get other text
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue
            
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    print(f"  - {paragraph.text.strip()}")

if __name__ == "__main__":
    pptx_path = r"c:\Users\LLOYD\Downloads\nasa data\assessment_deliverables\Assessment 1 Template.pptx"
    extract_pptx_text(pptx_path)
